"""
make_ppt.py
-----------
Generates the Surface Dipole Library presentation (.pptx)
Run:  python make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x2B, 0x55)   # dark navy  – slide backgrounds
TEAL       = RGBColor(0x00, 0x87, 0x8A)   # teal       – headings / accents
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY  = RGBColor(0x33, 0x33, 0x33)
ORANGE     = RGBColor(0xFF, 0x6B, 0x35)   # highlight / call-outs
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW     = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, size=16, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, space_before=Pt(4), italic=False):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def bg(slide, color=NAVY):
    """Fill entire slide background."""
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)


def accent_bar(slide, color=TEAL, h=0.07):
    add_rect(slide, 0, 0, 13.33, h, fill=color)


def slide_header(slide, title, subtitle=None):
    accent_bar(slide, TEAL)
    add_text(slide, title, 0.4, 0.12, 12, 0.55,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.65, 12, 0.4,
                 size=15, color=RGBColor(0xCC, 0xE5, 0xFF), italic=True)


def bullet_box(slide, items, l, t, w, h,
               title=None, title_color=TEAL,
               bullet="  •  ", size=15, bg_color=LIGHT_GREY,
               text_color=DARK_GREY, title_size=17):
    add_rect(slide, l, t, w, h, fill=bg_color)
    cur_t = t + 0.12
    if title:
        add_text(slide, title, l+0.15, cur_t, w-0.3, 0.38,
                 size=title_size, bold=True, color=title_color)
        cur_t += 0.38
    item_h = (h - 0.12 - (0.38 if title else 0)) / max(len(items), 1)
    item_h = max(item_h, 0.30)
    for item in items:
        add_text(slide, bullet + item, l+0.1, cur_t, w-0.2, item_h,
                 size=size, color=text_color)
        cur_t += item_h


def code_box(slide, code_text, l, t, w, h, font_size=11):
    add_rect(slide, l, t, w, h, fill=RGBColor(0x1E, 0x1E, 0x2E))
    add_text(slide, code_text, l+0.15, t+0.1, w-0.3, h-0.15,
             size=font_size, color=RGBColor(0xA6, 0xE2, 0x2E),
             bold=False, wrap=True)


def divider(slide, y, color=TEAL, thickness=0.025):
    add_rect(slide, 0.4, y, 12.5, thickness, fill=color)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, NAVY)

# gradient band
add_rect(sl, 0, 2.6, 13.33, 2.6, fill=RGBColor(0x06, 0x1A, 0x40))

add_text(sl, "Surface Dipole Library",
         0.8, 1.0, 11.7, 1.2,
         size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "A fast pre-screening tool for thin-film interface interactions",
         0.8, 2.1, 11.7, 0.55,
         size=20, color=RGBColor(0x7E, 0xC8, 0xE3),
         align=PP_ALIGN.CENTER, italic=True)

divider(sl, 2.75, TEAL, 0.04)

add_text(sl, "Calculates dipole moments & dipole–dipole interaction energies\n"
             "between surface atoms of a substrate and a thin film",
         1.5, 2.9, 10.3, 0.8,
         size=17, color=RGBColor(0xCC, 0xE5, 0xFF),
         align=PP_ALIGN.CENTER)

add_text(sl, "Shubham Gangwar   |   Department Project   |   2026",
         0.8, 6.7, 11.7, 0.5,
         size=14, color=RGBColor(0x88, 0x99, 0xAA),
         align=PP_ALIGN.CENTER)

add_text(sl, "v 1.0.0",
         11.5, 6.7, 1.5, 0.4,
         size=12, color=TEAL, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Presentation Agenda")

items = [
    "01   The Problem — Why this library was needed",
    "02   Library Architecture — How it was built",
    "03   What It Calculates — Key quantities",
    "04   Input Requirements — What you need to provide",
    "05   How to Get the Inputs — Coordinates & Charges",
    "06   Formulas & Calculation Flow — Step by step",
    "07   Output & Results — What the library returns",
    "08   Case Study — Borosilicate Glass vs. TiO₂",
    "09   Actionable Findings — What the results mean",
    "10   Conclusion & Future Work",
]

add_rect(sl, 0.5, 1.1, 12.3, 6.0, fill=LIGHT_GREY)
t = 1.25
for item in items:
    num, rest = item.split("   ", 1)
    add_text(sl, num, 0.7, t, 0.8, 0.47, size=15, bold=True, color=TEAL)
    add_text(sl, rest, 1.5, t, 11.0, 0.47, size=15, color=DARK_GREY)
    t += 0.53


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Problem
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "The Problem", "Why was this library needed?")

# left panel
add_rect(sl, 0.4, 1.1, 5.9, 5.8, fill=RGBColor(0xFF, 0xEB, 0xEB))
add_text(sl, "Without This Library", 0.6, 1.2, 5.5, 0.45,
         size=17, bold=True, color=RED)
problems = [
    "DFT calculation per config: 6–8 hours",
    "Screening 132 pairs = 3–5 days on HPC",
    "No guidance on which atoms to focus on",
    "Chemical intuition often wrong",
    "Expensive compute budget wasted on\nnon-dominant pairs",
]
t = 1.75
for p in problems:
    add_text(sl, "✗  " + p, 0.6, t, 5.5, 0.55, size=14, color=RED)
    t += 0.6

# right panel
add_rect(sl, 6.7, 1.1, 6.2, 5.8, fill=RGBColor(0xE8, 0xF8, 0xF0))
add_text(sl, "With This Library", 6.9, 1.2, 5.8, 0.45,
         size=17, bold=True, color=GREEN)
solutions = [
    "All 132 pairs evaluated: < 1 second",
    "Dominant pair identified instantly",
    "Atomic-level resolution — exact sites",
    "Physically grounded (force-field charges)",
    "Reduces DFT workload by 90–95%",
]
t = 1.75
for s in solutions:
    add_text(sl, "✓  " + s, 6.9, t, 5.8, 0.55, size=14, color=GREEN)
    t += 0.6

# centre vs arrow
add_text(sl, "VS", 6.0, 3.6, 0.7, 0.6,
         size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Library Architecture
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Library Architecture", "7-module Python package")

modules = [
    ("atoms.py",       "Atom dataclass\nwith labelling helpers",  "0.4"),
    ("dipole.py",      "calculate_dipole()\natomic_local_dipole()", "2.3"),
    ("geometry.py",    "distance()\npairwise_distances()",         "4.2"),
    ("interaction.py", "dipole_dipole_energy()\nSI-derived, eV",   "6.1"),
    ("ranking.py",     "rank_interactions()\nprint_summary()",      "8.0"),
    ("io.py",          "XYZ / CIF / POSCAR\nASE / Pymatgen readers","9.9"),
    ("main.py",        "analyze_surface_interactions()\nhigh-level API","11.8"),
]

for (name, desc, lx) in modules:
    lf = float(lx)
    add_rect(sl, lf, 1.2, 1.7, 1.3, fill=NAVY)
    add_text(sl, name, lf+0.05, 1.25, 1.6, 0.45,
             size=12, bold=True, color=TEAL)
    add_text(sl, desc, lf+0.05, 1.7, 1.6, 0.75,
             size=10, color=WHITE)

# flow arrow
add_text(sl, "User Input  →  io.py  →  atoms.py  →  dipole.py  →  geometry.py  →  interaction.py  →  ranking.py  →  main.py  →  Results",
         0.4, 2.8, 12.5, 0.45, size=12, color=TEAL, italic=True)

# __init__.py exposes all
add_rect(sl, 0.4, 3.35, 12.5, 0.55, fill=RGBColor(0x0D, 0x47, 0x6B))
add_text(sl, "__init__.py  —  re-exports everything via  from surface_dipole_library import ...",
         0.6, 3.4, 12.1, 0.45, size=13, bold=True, color=WHITE)

# install
add_rect(sl, 0.4, 4.05, 6.0, 1.1, fill=RGBColor(0x1E, 0x1E, 0x2E))
add_text(sl, "# Install\npip install surface-dipole-lib\n\n# Import\nfrom surface_dipole_library import analyze_surface_interactions",
         0.6, 4.1, 5.7, 1.0, size=11, color=RGBColor(0xA6, 0xE2, 0x2E))

add_rect(sl, 6.8, 4.05, 6.0, 1.1, fill=LIGHT_GREY)
add_text(sl, "Requirements", 7.0, 4.1, 5.7, 0.35,
         size=14, bold=True, color=NAVY)
add_text(sl, "  Python ≥ 3.8     numpy     pandas\n  ase  (optional, for CIF/POSCAR)\n  pymatgen  (optional, for oxidation states)",
         7.0, 4.45, 5.7, 0.65, size=12, color=DARK_GREY)

# pyproject note
add_text(sl, "pyproject.toml included — pip-installable package",
         0.4, 5.35, 12.5, 0.4, size=13, italic=True, color=TEAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What It Calculates
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "What the Library Calculates", "Three core quantities")

boxes = [
    ("1. Group Dipole Moment",
     "μ = Σᵢ qᵢ rᵢ\n\nA single vector summarising how polar\neach surface is overall.\n\nOutput: (μx, μy, μz) in e·Å and Debye",
     NAVY, 0.4),
    ("2. Local Atomic Dipole",
     "μᵢ = qᵢ (rᵢ − r̄)\n\nThe dipole contribution of each individual\natom relative to the group centroid.\n\nMakes calculation origin-independent.",
     RGBColor(0x0D, 0x47, 0x6B), 4.5),
    ("3. Dipole–Dipole\nInteraction Energy",
     "U = 14.3996 × [μ₁·μ₂/r³ − 3(μ₁·r)(μ₂·r)/r⁵]\n\nComputed for every substrate–film atom pair.\nRanked from most attractive to most repulsive.\n\nOutput: Energy in eV, labelled attractive/repulsive",
     RGBColor(0x00, 0x4D, 0x5A), 8.6),
]

for (title, body, col, lx) in boxes:
    add_rect(sl, lx, 1.1, 4.5, 5.8, fill=col)
    add_text(sl, title, lx+0.2, 1.2, 4.1, 0.7,
             size=16, bold=True, color=TEAL)
    divider(sl, 2.0 + (lx == 0.4)*0, TEAL, 0.02)
    add_text(sl, body, lx+0.2, 2.0, 4.1, 4.7,
             size=13, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Inputs Required
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Input Requirements", "What you must provide per atom")

# three cards
for i, (icon, label, detail, col) in enumerate([
    ("①", "Element Symbol",
     'e.g.  "Si",  "Ti",  "O",  "Al",  "Zn"\n\nIdentifies which atom it is.\nUsed for labelling and output tables.',
     NAVY),
    ("②", "3D Position  (Å)",
     'e.g.  [0.00, 0.00, 1.60]\n\nCartesian coordinates in Ångströms.\nBoth surfaces must share the same\ncoordinate system (same z-axis).',
     RGBColor(0x0D, 0x47, 0x6B)),
    ("③", "Partial Charge  (e)",
     'e.g.  +0.310  for Si,  −0.290  for O_NBO\n\nFraction of electron charge on the atom.\nNOT the formal oxidation state.\nSource: force field or DFT Bader analysis.',
     RGBColor(0x00, 0x4D, 0x5A)),
]):
    lx = 0.4 + i * 4.3
    add_rect(sl, lx, 1.1, 4.0, 4.5, fill=col)
    add_text(sl, icon, lx+0.15, 1.15, 0.8, 0.55,
             size=30, bold=True, color=TEAL)
    add_text(sl, label, lx+0.15, 1.7, 3.7, 0.5,
             size=17, bold=True, color=WHITE)
    add_text(sl, detail, lx+0.15, 2.3, 3.7, 3.2,
             size=13, color=RGBColor(0xCC, 0xE5, 0xFF))

# code example
add_rect(sl, 0.4, 5.75, 12.5, 1.45, fill=RGBColor(0x1E, 0x1E, 0x2E))
add_text(sl, 'substrate_atoms = [\n'
             '    {"element": "Si",  "charge": +0.310, "position": [0.00, 0.00, 0.00], "label": "Si1_sub"},\n'
             '    {"element": "O",   "charge": -0.290, "position": [0.00, 0.00, 1.60], "label": "O_NBO1_sub"},\n'
             ']',
         0.6, 5.82, 12.1, 1.3,
         size=11, color=RGBColor(0xA6, 0xE2, 0x2E))


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — How to Get Coordinates
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "How to Get Inputs — 3D Coordinates",
             "Three routes from fastest to most accurate")

routes = [
    ("Route A  — Download + VESTA",
     "1. Go to materialsproject.org\n"
     "2. Search material → Download CIF\n"
     "3. Open in VESTA → set Miller index\n"
     "4. Add vacuum layer → Export POSCAR/XYZ\n"
     "5. Load into library\n\nTime: ~10 minutes",
     GREEN, "Fastest\nGood for screening"),
    ("Route B  — ASE Slab Builder",
     "from ase.build import surface, bulk, add_vacuum\n"
     "tio2 = bulk('TiO2', crystalstructure='rutile',\n"
     "            a=4.594, c=2.958)\n"
     "slab = surface(tio2, (1,1,0), layers=4)\n"
     "add_vacuum(slab, 15.0)\n\nTime: ~20 minutes",
     YELLOW, "Reproducible\nScriptable"),
    ("Route C  — DFT Relaxation",
     "1. Build initial slab (Route A or B)\n"
     "2. Run VASP / Quantum ESPRESSO\n"
     "3. Use final CONTCAR / output.xyz\n"
     "   (true equilibrium positions)\n\n"
     "Time: hours–days on HPC",
     ORANGE, "Most accurate\nFor publications"),
]

for i, (title, body, col, tag) in enumerate(routes):
    lx = 0.4 + i * 4.3
    add_rect(sl, lx, 1.1, 4.0, 4.6, fill=LIGHT_GREY)
    add_rect(sl, lx, 1.1, 4.0, 0.45, fill=col)
    add_text(sl, title, lx+0.12, 1.12, 3.8, 0.4,
             size=14, bold=True, color=WHITE)
    add_text(sl, body, lx+0.12, 1.65, 3.8, 3.6,
             size=11, color=DARK_GREY)
    add_rect(sl, lx, 5.5, 4.0, 0.45, fill=col)
    add_text(sl, tag, lx+0.12, 5.52, 3.8, 0.4,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# filter tip
add_rect(sl, 0.4, 6.1, 12.5, 0.55, fill=RGBColor(0xEA, 0xF4, 0xFF))
add_text(sl, "Tip — Keep only surface layer atoms:  "
             "surface_atoms = [a for a in all_atoms if a.position[2] >= z_max − 3.0]",
         0.6, 6.15, 12.1, 0.45, size=12, color=NAVY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — How to Get Partial Charges
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "How to Get Inputs — Partial Charges",
             "Charge ≠ Oxidation State")

# explanation box
add_rect(sl, 0.4, 1.1, 12.5, 1.2, fill=RGBColor(0xFF, 0xF3, 0xCD))
add_text(sl, "Why not use formal oxidation states?  Si is +4 formally, but using +4 makes forces 167× too strong "
             "— simulations collapse.\n"
             "Partial charges reflect the real ~50% covalent nature of Si–O bonds.  "
             "Si = +0.310 e because it retains 3.69 of 4 valence electrons.",
         0.6, 1.15, 12.1, 1.1, size=13, color=DARK_GREY)

# three source cards
srcs = [
    ("Published Force Fields\n(Fastest, Good accuracy)",
     "ClayFF  →  SiO₂, glass\n  Si: +0.310,  O_BO: −0.155,  O_NBO: −0.290\n\n"
     "Matsui–Akaogi  →  TiO₂\n  Ti5c: +0.580,  O: −0.290\n\n"
     "INTERFACE FF  →  Al₂O₃\n  Al: +0.480,  O: −0.320\n\n"
     "Look up: '[material] partial charge force field'",
     GREEN),
    ("DFT Bader Analysis\n(Most accurate)",
     "1. Run VASP with LAECHG=.TRUE.\n"
     "2. bader CHGCAR -ref AECCAR0 -ref AECCAR2\n"
     "3. Read ACF.dat  (column 5 = Bader charge)\n\n"
     "Partial charge = valence electrons − Bader charge\n\n"
     "Quantum ESPRESSO: projwfc.x → Löwdin charges\n"
     "CP2K: MULLIKEN_POP keyword",
     ORANGE),
    ("Pymatgen BVAnalyzer\n(Rough, for quick screening)",
     "from pymatgen.core import Structure\n"
     "from pymatgen.analysis.bond_valence\n"
     "    import BVAnalyzer\n\n"
     "s = Structure.from_file('Al2O3.cif')\n"
     "bv = BVAnalyzer()\n"
     "s = bv.get_oxi_state_decorated_structure(s)\n\n"
     "Gives integer oxidation states (+3, −2)\nNot partial — use as last resort",
     TEAL),
]

for i, (title, body, col) in enumerate(srcs):
    lx = 0.4 + i * 4.3
    add_rect(sl, lx, 2.45, 4.0, 4.6, fill=LIGHT_GREY)
    add_rect(sl, lx, 2.45, 4.0, 0.55, fill=col)
    add_text(sl, title, lx+0.12, 2.47, 3.8, 0.52,
             size=13, bold=True, color=WHITE)
    add_text(sl, body, lx+0.12, 3.1, 3.8, 3.85,
             size=10, color=DARK_GREY)

add_rect(sl, 0.4, 7.1, 12.5, 0.25, fill=TEAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Formulas & Calculation Flow
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Formulas & Calculation Flow", "Step-by-step inside the library")

steps = [
    ("STEP 1", "Group Dipole Moment",
     "μ = Σᵢ qᵢ rᵢ",
     "Sum of (charge × position) over all atoms.\nOutputs one vector per surface.\nUsed as diagnostic — tells you if surfaces are polar.",
     NAVY),
    ("STEP 2", "Geometric Centroid",
     "r̄ = (1/N) Σᵢ rᵢ",
     "Unweighted mean position of all atoms.\nUsed as origin for local dipoles.\nMakes all energies origin-independent.",
     RGBColor(0x0D, 0x47, 0x6B)),
    ("STEP 3", "Local Atomic Dipole",
     "μᵢ = qᵢ (rᵢ − r̄)",
     "Per-atom dipole relative to centroid.\nCaptures each atom's individual contribution.\nLarger charge + farther from centroid = stronger dipole.",
     RGBColor(0x00, 0x5F, 0x73)),
    ("STEP 4", "Pairwise Interaction Energy",
     "U = 14.3996 × [μ₁·μ₂/r³\n    − 3(μ₁·r̂)(μ₂·r̂)/r³]",
     "Applied to all N×M substrate–film pairs.\n14.3996 = Coulomb constant in eV·Å/e².\nU < 0 → attractive,  U > 0 → repulsive.",
     RGBColor(0x00, 0x4D, 0x5A)),
    ("STEP 5", "Rank & Filter",
     "Sort by U (ascending)\nApply distance cutoff",
     "132 pairs sorted from most negative to most positive.\nCutoff removes distant pairs (default 8 Å).\nOutputs ranked table with labels and types.",
     RGBColor(0x00, 0x3B, 0x45)),
]

lx = 0.3
for (tag, title, formula, desc, col) in steps:
    add_rect(sl, lx, 1.1, 2.3, 5.9, fill=col)
    add_text(sl, tag,   lx+0.1, 1.15, 2.1, 0.4, size=12, bold=True, color=TEAL)
    add_text(sl, title, lx+0.1, 1.55, 2.1, 0.45, size=13, bold=True, color=WHITE)
    add_rect(sl, lx+0.1, 2.1, 2.1, 0.9, fill=RGBColor(0x05, 0x10, 0x25))
    add_text(sl, formula, lx+0.15, 2.12, 2.0, 0.86,
             size=11, bold=True, color=TEAL)
    add_text(sl, desc, lx+0.1, 3.1, 2.1, 3.7,
             size=11, color=RGBColor(0xCC, 0xE5, 0xFF))
    lx += 2.55


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Why 0.310 for Si
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Deep Dive — Why Si = +0.310 e and Not +4",
             "Partial charge vs formal oxidation state")

# left: explanation
add_rect(sl, 0.4, 1.1, 6.2, 5.9, fill=LIGHT_GREY)
add_text(sl, "The Si–O Bond Is ~50% Covalent", 0.6, 1.2, 5.8, 0.45,
         size=16, bold=True, color=NAVY)
explanation = (
    "Formal oxidation state:  Si⁴⁺  →  +4\n"
    "This assumes 4 electrons fully transferred to O.\n\n"
    "Reality: Si–O is partly covalent. Electrons are\n"
    "shared, not transferred. Si retains 3.69 of 4\n"
    "valence electrons → net charge = +0.310 e only.\n\n"
    "Using +4 would make forces 167× too strong.\n\n"
    "The 0.310 value comes from ClayFF force field\n"
    "(Cygan et al., J. Phys. Chem. B, 2004):\n"
    "  → Quantum calc on SiO₄(OH)₄ cluster\n"
    "  → ESP fitting to electron density\n"
    "  → Charge balanced across 4 bonds:\n"
    "       Si donates 0.31 e total\n"
    "       Each BO gets 0.31/4 × 2 = 0.155 e\n"
    "       → O_BO charge = −0.155 e  ✓\n"
    "       NBO gets more → −0.290 e  ✓"
)
add_text(sl, explanation, 0.6, 1.75, 5.8, 5.1, size=12, color=DARK_GREY)

# right: electronegativity table + diagram
add_rect(sl, 6.9, 1.1, 6.0, 2.7, fill=NAVY)
add_text(sl, "Electronegativity Drives the Charge", 7.1, 1.15, 5.7, 0.45,
         size=14, bold=True, color=TEAL)
en_text = (
    "Bond     ΔEN    Partial charge on metal\n"
    "──────────────────────────────────────\n"
    "Ti–O     1.90    Ti = +0.580 e  (highest)\n"
    "Si–O     1.54    Si = +0.310 e\n"
    "B–O      1.40    B  = +0.270 e\n"
    "H–O      1.24    H  = +0.175 e  (lowest)\n\n"
    "Larger ΔEN → more electrons pulled away\n"
    "→ higher positive partial charge"
)
add_text(sl, en_text, 7.1, 1.65, 5.7, 2.0, size=11, color=WHITE)

add_rect(sl, 6.9, 3.95, 6.0, 3.05, fill=LIGHT_GREY)
add_text(sl, "Charge Balance Check (SiO₂ unit cell)", 7.1, 4.0, 5.7, 0.4,
         size=14, bold=True, color=NAVY)
balance = (
    "Si charge:  +0.310 e × 1 = +0.310 e\n"
    "O  charge:  −0.155 e × 2 = −0.310 e\n"
    "─────────────────────────────────────\n"
    "Net charge:             =   0.000 e  ✓\n\n"
    "The unit cell is electrically neutral.\n"
    "Charges donated by Si exactly equal\n"
    "charges received by the two oxygens."
)
add_text(sl, balance, 7.1, 4.5, 5.7, 2.4, size=12, color=DARK_GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Case Study Setup
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Case Study — Borosilicate Glass vs. TiO₂ Thin Film",
             "Real-world interface test")

# interface diagram (text-based)
add_rect(sl, 0.4, 1.1, 5.6, 5.9, fill=NAVY)
add_text(sl, "Interface Geometry  (z-axis view)", 0.6, 1.15, 5.2, 0.4,
         size=14, bold=True, color=TEAL)
diagram = (
    "z (Å)\n"
    " 5.8 ─  H_OH1_film   (+0.18 e)\n"
    " 5.2 ─  O_OH1_film   (−0.20 e)\n"
    " 4.5 ─  Ti6c1_film   (+0.48 e)\n"
    " 4.2 ─  O_br_film    (−0.29 e)\n"
    " 3.8 ─  O_ip_film    (−0.26 e)\n"
    " 3.5 ─  Ti5c_film  ★ (+0.58 e)  ←── primary\n"
    "                                         anchor\n"
    "  ══════ interface gap ≈ 1.9 Å ══════\n\n"
    " 1.6 ─  O_NBO_sub  ★ (−0.29 e)  ←── primary\n"
    "                                         docking\n"
    " 0.8 ─  O_BO_sub     (−0.155 e)\n"
    " 0.0 ─  Si_sub / B   (+0.31 e)\n\n"
    "  ↑ Film dipole (+z)   ↓ Glass dipole (−z)\n"
    "  Anti-parallel → macroscopic attraction"
)
add_text(sl, diagram, 0.6, 1.65, 5.2, 5.2, size=10, color=WHITE)

# right: atom inventory
add_rect(sl, 6.2, 1.1, 6.7, 2.8, fill=LIGHT_GREY)
add_text(sl, "Glass Substrate — 12 atoms", 6.4, 1.15, 6.3, 0.4,
         size=14, bold=True, color=NAVY)
glass_info = (
    "Si × 3    +0.310 e    Tetrahedral centres\n"
    "O_BO × 4  −0.155 e   Bridging oxygens\n"
    "O_NBO × 3 −0.290 e   Non-bridging (surface)\n"
    "B × 1     +0.270 e   Borosilicate network\n"
    "H × 1     +0.175 e   Silanol –OH"
)
add_text(sl, glass_info, 6.4, 1.65, 6.3, 2.1, size=12, color=DARK_GREY)

add_rect(sl, 6.2, 4.05, 6.7, 2.7, fill=LIGHT_GREY)
add_text(sl, "TiO₂ Film — 11 atoms", 6.4, 4.1, 6.3, 0.4,
         size=14, bold=True, color=NAVY)
film_info = (
    "Ti5c × 3  +0.580 e   5-fold Ti (reactive)\n"
    "Ti6c × 1  +0.480 e   6-fold Ti (bulk-like)\n"
    "O_br × 3  −0.290 e   Bridging O rows\n"
    "O_ip × 2  −0.260 e   In-plane O\n"
    "O_OH × 1  −0.200 e   Hydroxyl O\n"
    "H_OH × 1  +0.180 e   Hydroxyl H"
)
add_text(sl, film_info, 6.4, 4.55, 6.3, 2.1, size=12, color=DARK_GREY)

add_rect(sl, 6.2, 6.85, 6.7, 0.45, fill=TEAL)
add_text(sl, "Total pairs to evaluate:  12 × 11 = 132",
         6.4, 6.88, 6.3, 0.38, size=14, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Test Results: Dipole Moments & Unit Tests
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Case Study Results — Dipole Moments & Unit Tests")

# unit tests
add_rect(sl, 0.4, 1.1, 6.0, 5.9, fill=LIGHT_GREY)
add_text(sl, "Unit Test Results — All 6 Passed", 0.6, 1.15, 5.6, 0.45,
         size=15, bold=True, color=GREEN)
tests = [
    ("test_dipole_moments",              "Glass: 7.748 D  |  TiO₂: 13.463 D"),
    ("test_interaction_table_structure", "132 pairs generated correctly"),
    ("test_best_interaction_attractive", "Best = O_NBO1 – Ti5c1 at −2.224 eV"),
    ("test_distance_calculation",        "d = 1.9000 Å verified"),
    ("test_cutoff_filter",               "132 → 65 pairs at 4.5 Å cutoff"),
    ("test_energy_symmetry",             "U(A→B) = U(B→A) confirmed"),
]
t = 1.7
for (name, detail) in tests:
    add_rect(sl, 0.5, t, 5.8, 0.75, fill=RGBColor(0xE8, 0xF8, 0xF0))
    add_text(sl, "✓", 0.6, t+0.05, 0.4, 0.55, size=16, bold=True, color=GREEN)
    add_text(sl, name, 1.05, t+0.05, 4.7, 0.35, size=12, bold=True, color=DARK_GREY)
    add_text(sl, detail, 1.05, t+0.38, 4.7, 0.3, size=11, color=TEAL, italic=True)
    t += 0.83

# dipole moments
add_rect(sl, 6.7, 1.1, 6.3, 2.8, fill=NAVY)
add_text(sl, "Glass Dipole Moment", 6.9, 1.15, 5.9, 0.4,
         size=14, bold=True, color=TEAL)
add_text(sl,
         "μx = −0.6184 e·Å\n"
         "μy = −0.1971 e·Å\n"
         "μz = −1.4768 e·Å   ← dominant\n"
         "|μ| = 1.6131 e·Å  =  7.748 Debye\n"
         "Net Q = −0.115 e\n\n"
         "Dipole points DOWN (−z)\n"
         "Negative charge at top of glass (NBO layer)",
         6.9, 1.65, 5.9, 2.1, size=11, color=WHITE)

add_rect(sl, 6.7, 4.05, 6.3, 2.8, fill=RGBColor(0x0D, 0x47, 0x6B))
add_text(sl, "TiO₂ Film Dipole Moment", 6.9, 4.1, 5.9, 0.4,
         size=14, bold=True, color=TEAL)
add_text(sl,
         "μx = −0.0592 e·Å\n"
         "μy = +0.9833 e·Å\n"
         "μz = +2.6240 e·Å   ← dominant\n"
         "|μ| = 2.8028 e·Å  =  13.463 Debye\n"
         "Net Q = +0.810 e\n\n"
         "Dipole points UP (+z)\n"
         "Positive Ti5c layer faces the glass",
         6.9, 4.55, 5.9, 2.1, size=11, color=WHITE)

add_rect(sl, 6.7, 6.95, 6.3, 0.35, fill=TEAL)
add_text(sl, "Anti-parallel dipoles → macroscopic attraction confirmed",
         6.9, 6.97, 5.9, 0.3, size=12, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Ranked Table (Top 10 + Bottom 5)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Case Study Results — Ranked Interaction Table",
             "132 pairs sorted by energy — top 10 attractive + bottom 5 repulsive")

# header row
cols_w  = [0.55, 2.3, 2.3, 1.7, 1.9, 1.8]
cols_x  = [0.3, 0.88, 3.22, 5.55, 7.28, 9.21]
headers = ["Rank", "Substrate Atom", "Film Atom", "Distance (Å)", "Energy (eV)", "Type"]

add_rect(sl, 0.3, 1.1, 11.15, 0.42, fill=NAVY)
for hdr, cx, cw in zip(headers, cols_x, cols_w):
    add_text(sl, hdr, cx, 1.12, cw, 0.37, size=12, bold=True, color=WHITE)

rows = [
    ("1",  "O_NBO1_sub", "Ti5c1_film", "1.9000", "−2.2243", "attractive", True),
    ("2",  "O_NBO3_sub", "Ti5c3_film", "2.4473", "−0.3509", "attractive", False),
    ("3",  "H_OH1_sub",  "Ti5c1_film", "1.4980", "−0.1875", "attractive", True),
    ("4",  "H_OH1_sub",  "O_br1_film", "1.9217", "−0.1707", "attractive", False),
    ("5",  "Si2_sub",    "O_br2_film", "4.2455", "−0.1227", "attractive", True),
    ("6",  "O_NBO1_sub", "Ti5c3_film", "4.0451", "−0.1217", "attractive", False),
    ("7",  "O_BO3_sub",  "Ti5c3_film", "2.9082", "−0.1162", "attractive", True),
    ("8",  "O_BO1_sub",  "Ti5c1_film", "3.1385", "−0.1040", "attractive", False),
    ("9",  "Si2_sub",    "Ti5c2_film", "4.0817", "−0.1000", "attractive", True),
    ("10", "O_BO2_sub",  "Ti5c2_film", "2.8276", "−0.0996", "attractive", False),
]

t = 1.55
for (rank, sub, film, dist, energy, typ, alt) in rows:
    bg_c = RGBColor(0xE8,0xF8,0xF0) if alt else WHITE
    if rank == "1":
        bg_c = RGBColor(0xD4,0xEF,0xDF)
    add_rect(sl, 0.3, t, 11.15, 0.38, fill=bg_c)
    vals = [rank, sub, film, dist, energy, typ]
    for val, cx, cw in zip(vals, cols_x, cols_w):
        col = GREEN if typ == "attractive" else RED
        if val == energy:
            tc = GREEN if "−" in energy else RED
        elif val == rank and rank == "1":
            tc = NAVY
        else:
            tc = DARK_GREY
        add_text(sl, val, cx, t+0.02, cw, 0.34, size=11,
                 bold=(rank=="1"), color=tc)
    t += 0.39

# separator
add_rect(sl, 0.3, t+0.04, 11.15, 0.04, fill=ORANGE)
t += 0.14

# bottom 5 repulsive
rep_rows = [
    ("128", "O_NBO1_sub", "O_ip1_film",  "2.7321", "+0.1594", "repulsive", True),
    ("129", "O_NBO2_sub", "Ti5c2_film",  "2.8320", "+0.2178", "repulsive", False),
    ("130", "Si1_sub",    "Ti5c1_film",  "3.5000", "+0.2451", "repulsive", True),
    ("131", "O_NBO3_sub", "O_br3_film",  "2.8674", "+0.2996", "repulsive", False),
    ("132", "O_NBO2_sub", "O_br2_film",  "2.6729", "+0.5704", "repulsive", True),
]

for (rank, sub, film, dist, energy, typ, alt) in rep_rows:
    bg_c = RGBColor(0xFF,0xEB,0xEB) if alt else RGBColor(0xFF,0xF5,0xF5)
    if rank == "132":
        bg_c = RGBColor(0xF5,0xCC,0xCC)
    add_rect(sl, 0.3, t, 11.15, 0.38, fill=bg_c)
    vals = [rank, sub, film, dist, energy, typ]
    for val, cx, cw in zip(vals, cols_x, cols_w):
        tc = RED if typ == "repulsive" else GREEN
        if val not in (energy, typ):
            tc = DARK_GREY
        add_text(sl, val, cx, t+0.02, cw, 0.34, size=11,
                 bold=(rank=="132"), color=tc)
    t += 0.39

# stats bar
add_rect(sl, 0.3, 7.15, 11.15, 0.3, fill=NAVY)
add_text(sl, "Total: 132 pairs   |   Attractive: 67 (50.8%)   |   "
             "Repulsive: 65 (49.2%)   |   Best: −2.2243 eV   |   Worst: +0.5704 eV",
         0.45, 7.17, 10.8, 0.26, size=11, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Actionable Findings
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Actionable Findings from the Case Study",
             "What the ranked table tells the experimentalist")

findings = [
    ("Finding 1",
     "Dominant pair: O_NBO ↔ Ti5c  at  −2.224 eV",
     "6.3× stronger than rank 2. Focus all surface engineering here.\n"
     "Action: Increase NBO density on glass by reducing B₂O₃ content\nor alkali leaching before deposition.",
     GREEN),
    ("Finding 2",
     "Ti5c is 50× more active than Ti6c",
     "Ti5c ranks 1,2,3,6,7,8,9,10.  Ti6c first appears at rank 21 (−0.045 eV).\n"
     "Action: Deposit TiO₂ under reduced O₂ pressure to maximise Ti5c surface density.",
     TEAL),
    ("Finding 3",
     "Surface –OH groups HELP adhesion",
     "Silanol H contributes ranks 3 & 4 (−0.188 + −0.171 = −0.36 eV extra).\n"
     "Action: Wet-clean glass before deposition. Do not avoid hydroxylation.",
     ORANGE),
    ("Finding 4",
     "O–O misregistry is the worst case  (+0.570 eV)",
     "O_NBO on glass directly above O_br on film → strongly repulsive.\n"
     "Action: Optimise in-plane rotation of film during deposition.",
     RED),
    ("Finding 5",
     "Library reduces DFT workload by 90–95%",
     "132 pairs in <1 sec. Cutoff at 3 Å leaves ~10 high-priority DFT targets.\n"
     "Action: Run DFT only on top-ranked pairs identified here.",
     NAVY),
]

t = 1.15
for (tag, title, body, col) in findings:
    add_rect(sl, 0.4, t, 12.5, 1.08, fill=LIGHT_GREY)
    add_rect(sl, 0.4, t, 1.5, 1.08, fill=col)
    add_text(sl, tag, 0.42, t+0.05, 1.45, 0.4,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 2.05, t+0.05, 10.7, 0.38,
             size=14, bold=True, color=col)
    add_text(sl, body, 2.05, t+0.45, 10.7, 0.55,
             size=11, color=DARK_GREY)
    t += 1.17


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Limitations & Future Work
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, WHITE)
slide_header(sl, "Limitations & Future Work")

add_rect(sl, 0.4, 1.1, 6.0, 5.9, fill=RGBColor(0xFF, 0xEB, 0xEB))
add_text(sl, "Current Limitations", 0.6, 1.15, 5.6, 0.45,
         size=16, bold=True, color=RED)
lims = [
    "Fixed partial charges — no polarizability.\nInduction energy (~15–20% of rank-1 pair)\nnot captured.",
    "No dispersion / van der Waals (London C₆/r⁶)\nterm included.",
    "No charge–charge Coulomb term.\nElectrostatic monopole–monopole interaction\nnot calculated.",
    "Static model — no thermal motion,\nno solvent screening, no dynamics.",
    "Accuracy depends on quality of input charges.\nForce-field charges ≈ good; formal oxidation\nstates = rough.",
]
t = 1.7
for lim in lims:
    add_text(sl, "⚠  " + lim, 0.6, t, 5.5, 0.9, size=12, color=DARK_GREY)
    t += 0.98

add_rect(sl, 6.7, 1.1, 6.3, 5.9, fill=RGBColor(0xE8, 0xF8, 0xF0))
add_text(sl, "Planned Future Work", 6.9, 1.15, 5.9, 0.45,
         size=16, bold=True, color=GREEN)
futures = [
    "Add polarizability (α) per element →\nincludes induction energy term.",
    "Integrate London dispersion (C₆/r⁶)\nfor van der Waals contribution.",
    "GUI / web interface for non-programmers.",
    "Extend to liquid–solid interfaces\n(solvent screening factor ε).",
    "Benchmark against experimental adhesion\nenergy database for validation.",
    "Auto-suggest deposition conditions based\non ranked output.",
]
t = 1.7
for fut in futures:
    add_text(sl, "→  " + fut, 6.9, t, 5.9, 0.9, size=12, color=DARK_GREY)
    t += 0.98


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Conclusion
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl, NAVY)

add_text(sl, "Conclusion", 0.8, 0.3, 11.7, 0.75,
         size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
divider(sl, 1.1, TEAL, 0.04)

points = [
    ("What was built",
     "A production-ready Python library — 7 modules, pip-installable — that screens\n"
     "all substrate–film atom-pair interactions using dipole–dipole physics."),
    ("What it calculates",
     "Group dipole moments, per-atom local dipoles, and pairwise interaction\n"
     "energies for all N×M atom combinations, ranked from most attractive to most repulsive."),
    ("What the Glass–TiO₂ test proved",
     "O_NBO ↔ Ti5c is the dominant bonding pair at −2.224 eV — 6.3× stronger than\n"
     "rank 2. All 6 unit tests pass. Results match DFT literature expectations."),
    ("Why it matters",
     "Reduces interface screening from 3–5 days of DFT to under 1 second.\n"
     "Gives experimentalists specific, actionable guidance on surface engineering."),
]

t = 1.3
for (title, body) in points:
    add_rect(sl, 0.6, t, 12.1, 1.3, fill=RGBColor(0x06, 0x1A, 0x40))
    add_rect(sl, 0.6, t, 0.08, 1.3, fill=TEAL)
    add_text(sl, title, 0.85, t+0.08, 11.6, 0.38,
             size=15, bold=True, color=TEAL)
    add_text(sl, body,  0.85, t+0.46, 11.6, 0.78,
             size=13, color=WHITE)
    t += 1.42

add_text(sl, "github.com/shubhamgangwar-01/Surface-Dipole-Library",
         0.8, 7.1, 11.7, 0.35,
         size=13, color=TEAL, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
out = "Surface_Dipole_Library_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
